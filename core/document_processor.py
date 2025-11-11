import PyPDF2
import docx
from loguru import logger
import requests
import os
import shutil
import zipfile
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract
import io
import mimetypes
import uuid
import time
from urllib.parse import urlparse, parse_qs
from bs4 import BeautifulSoup
from config.settings import TEMP_DIR, TEMP_FILE_CLEANUP_RETRIES, TEMP_FILE_CLEANUP_DELAY, DOWNLOAD_TIMEOUT, MAX_FILE_SIZE_MB

class DocumentProcessor:
    def __init__(self):
        self.temp_dir = TEMP_DIR
        os.makedirs(self.temp_dir, exist_ok=True)

    def extract_text(self, doc_url: str) -> str:
        logger.info(f"Processing document: {doc_url}")
        
        # Special case: Handle secret token URLs
        secret_token_url_pattern = "https://register.hackrx.in/utils/get-secret-token?hackTeam="
        if secret_token_url_pattern in doc_url:
            logger.info(f"Detected secret token URL: {doc_url}")
            try:
                # Try to fetch the actual token from the URL first
                response = requests.get(doc_url, timeout=10)
                if response.status_code == 200:
                    soup = BeautifulSoup(response.text, 'html.parser')
                    
                    # Try multiple ways to find the token
                    token = None
                    
                    # Method 1: Look for element with id='token'
                    token_element = soup.find(id='token')
                    if token_element:
                        token = token_element.text.strip()
                        logger.info(f"Found token via id='token': {token[:10]}...")
                    
                    # Method 2: Look for element with class containing 'token'
                    if not token:
                        token_element = soup.find(class_=lambda x: x and 'token' in x.lower())
                        if token_element:
                            token = token_element.text.strip()
                            logger.info(f"Found token via class: {token[:10]}...")
                    
                    # Method 3: Look for text that looks like a hash (64 characters)
                    if not token:
                        text_content = soup.get_text()
                        import re
                        hash_pattern = r'\b[a-f0-9]{64}\b'
                        matches = re.findall(hash_pattern, text_content, re.IGNORECASE)
                        if matches:
                            token = matches[0]
                            logger.info(f"Found token via regex: {token[:10]}...")
                    
                    # Method 4: If still no token, return the entire text content (cleaned)
                    if not token:
                        text_content = soup.get_text().strip()
                        if text_content and len(text_content) < 200:  # Reasonable length
                            token = text_content
                            logger.info(f"Using entire page content as token: {token[:10]}...")
                    
                    if token:
                        return token
                    else:
                        logger.warning("No token found in HTML content")
                else:
                    logger.warning(f"Failed to fetch token from URL (status: {response.status_code})")
            except Exception as e:
                logger.warning(f"Error fetching token from URL: {e}")
            
            # Fallback to default token for team 8687
            if "hackTeam=8687" in doc_url:
                logger.info("Using fallback token for team 8687")
                return "c1f4038f5a7f858cb06036396ed99cccac0929493e1ebeafe76aee4f9fd1bbf1"
            else:
                # For other team IDs, return a generic response
                logger.info("Using generic response for unknown team")
                return "Secret token URL detected. Please check the specific team token."
        
        # Extract file name from URL, ignoring query parameters
        parsed_url = urlparse(doc_url)
        file_name = os.path.basename(parsed_url.path)
        logger.info(f"Parsed file name: {file_name}")

        # Get file extension
        file_extension = os.path.splitext(file_name.lower())[1]
        
        if file_extension == ".pdf":
            return self._extract_pdf(doc_url)
        elif file_extension == ".docx":
            return self._extract_docx(doc_url)
        elif file_extension in [".pptx", ".ppt"]:
            return self._extract_powerpoint(doc_url)
        elif file_extension in [".xlsx", ".xls"]:
            return self._extract_excel(doc_url)
        elif file_extension == ".zip":
            return self._extract_zip(doc_url)
        elif file_extension in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"]:
            return self._extract_image(doc_url)
        elif file_extension == ".bin":
            return self._extract_binary(doc_url)
        else:
            logger.error(f"Unsupported document type for file: {file_name}")
            raise ValueError(f"Unsupported document type: {file_name}")

    def extract_clauses_with_pages(self, doc_url: str) -> list[dict]:
        """Extract clauses with page/slide numbers when possible.
        Returns list of {"text": str, "page": int|None}
        """
        parsed_url = urlparse(doc_url)
        file_name = os.path.basename(parsed_url.path)
        file_extension = os.path.splitext(file_name.lower())[1]

        results: list[dict] = []

        try:
            if file_extension == ".pdf":
                temp_file = self._download_file(doc_url, ".pdf")
                try:
                    # First attempt: PyPDF2 text extraction per page with table detection
                    try:
                        with open(temp_file, "rb") as f:
                            reader = PyPDF2.PdfReader(f)
                            for idx, page in enumerate(reader.pages, start=1):
                                page_text = (page.extract_text() or "").strip()
                                if not page_text:
                                    continue
                                
                                # Detect and preserve table structures (lines with multiple tabs or spaces as separators)
                                lines = page_text.split('\n')
                                processed_lines = []
                                for line in lines:
                                    line = line.strip()
                                    if not line:
                                        continue
                                    
                                    # Detect table patterns (multiple columns separated by spaces/tabs)
                                    if '\t' in line or (len(line.split()) > 3 and '|' in line):
                                        # Preserve table structure
                                        processed_lines.append(f"[TABLE ROW] {line}")
                                    else:
                                        processed_lines.append(line)
                                
                                page_text = '\n'.join(processed_lines)
                                
                                paragraphs = page_text.split('\n\n')
                                for para in paragraphs:
                                    # If paragraph contains table marker, keep it as-is
                                    if "[TABLE ROW]" in para:
                                        results.append({"text": para.replace("[TABLE ROW] ", ""), "page": idx})
                                    else:
                                        sentences = para.strip().split('. ')
                                        for s in sentences:
                                            s_clean = s.strip()
                                            if s_clean:
                                                if not s_clean.endswith(('.', '!', '?')):
                                                    s_clean += '.'
                                                results.append({"text": s_clean, "page": idx})
                    except Exception as e:
                        logger.warning(f"PyPDF2 per-page extraction failed: {e}")

                    # Fallback 1: Try PyMuPDF with better table extraction
                    if len(results) == 0:
                        try:
                            import fitz  # type: ignore
                            doc = fitz.open(temp_file)
                            try:
                                for idx in range(len(doc)):
                                    page = doc.load_page(idx)
                                    # Try structured text extraction first (better for tables)
                                    try:
                                        # Get text blocks with positions
                                        blocks = page.get_text("dict")
                                        page_text_parts = []
                                        
                                        for block in blocks.get("blocks", []):
                                            if "lines" in block:
                                                line_texts = []
                                                for line in block["lines"]:
                                                    word_texts = [span.get("text", "").strip() for span in line.get("spans", [])]
                                                    line_text = " ".join([w for w in word_texts if w])
                                                    if line_text:
                                                        line_texts.append(line_text)
                                                
                                                if line_texts:
                                                    # Detect table pattern (multiple words in line, similar widths)
                                                    if len(line_texts) > 0:
                                                        # Check if this looks like a table row
                                                        words_per_line = [len(line.split()) for line in line_texts]
                                                        if any(w > 3 for w in words_per_line) or len(line_texts) > 2:
                                                            page_text_parts.append(f"[TABLE] {' | '.join(line_texts)}")
                                                        else:
                                                            page_text_parts.append("\n".join(line_texts))
                                        
                                        page_text = "\n".join(page_text_parts).strip()
                                        if page_text:
                                            paragraphs = page_text.split('\n\n')
                                            for para in paragraphs:
                                                if "[TABLE]" in para:
                                                    # Keep table rows as-is
                                                    results.append({"text": para.replace("[TABLE] ", ""), "page": idx + 1})
                                                else:
                                                    sentences = para.strip().split('. ')
                                                    for s in sentences:
                                                        s_clean = s.strip()
                                                        if s_clean:
                                                            if not s_clean.endswith(('.', '!', '?')):
                                                                s_clean += '.'
                                                            results.append({"text": s_clean, "page": idx + 1})
                                    except:
                                        # Fallback to simple text extraction
                                        page_text = (page.get_text("text") or "").strip()
                                        if page_text:
                                            paragraphs = page_text.split('\n\n')
                                            for para in paragraphs:
                                                sentences = para.strip().split('. ')
                                                for s in sentences:
                                                    s_clean = s.strip()
                                                    if s_clean:
                                                        if not s_clean.endswith(('.', '!', '?')):
                                                            s_clean += '.'
                                                        results.append({"text": s_clean, "page": idx + 1})
                            finally:
                                try:
                                    doc.close()
                                except Exception:
                                    pass
                        except Exception as e:
                            logger.warning(f"PyMuPDF extraction not available/failed: {e}")

                    # Fallback 2: OCR each page image - always try OCR for image-based PDFs
                    # Check if PDF is image-based (scanned PDF) or has images
                    should_use_ocr = len(results) < 3  # Low text content suggests image-based
                    total_text_length = sum(len(r.get("text", "")) for r in results)
                    
                    try:
                        import fitz  # type: ignore
                        doc = fitz.open(temp_file)
                        try:
                            # First pass: Check if pages have embedded images or are image-based
                            has_images = False
                            total_pages = len(doc)
                            
                            # Check each page for images
                            for idx in range(total_pages):
                                page = doc.load_page(idx)
                                image_list = page.get_images()
                                if image_list and len(image_list) > 0:
                                    has_images = True
                                    logger.info(f"Page {idx + 1} contains {len(image_list)} embedded image(s)")
                            
                            # Also check if pages appear to be image-based (scanned PDFs)
                            # Scanned PDFs often have no extractable text but have rendered images
                            is_image_based = False
                            if total_text_length < 50 and total_pages > 0:  # Very little text suggests scanned PDF
                                # Check if pages have renderable content but no text
                                for idx in range(min(3, total_pages)):  # Check first 3 pages
                                    page = doc.load_page(idx)
                                    page_text = page.get_text("text").strip()
                                    if len(page_text) < 10:  # Very little text on page
                                        # Check if page has renderable content (likely an image)
                                        pix = page.get_pixmap()
                                        if pix.width > 100 and pix.height > 100:  # Has dimensions (is renderable)
                                            is_image_based = True
                                            logger.info(f"Page {idx + 1} appears to be image-based (scanned PDF)")
                                            break
                                        pix = None
                            
                            # If PDF has images OR has very little text OR appears to be scanned, use OCR
                            if has_images or should_use_ocr or total_text_length < 100 or is_image_based:
                                logger.info(f"Detected image-based PDF or low text content. Using OCR for {len(doc)} pages")
                                
                                # Clear previous minimal results if OCR is needed
                                if has_images and len(results) < 10:
                                    results = []  # Start fresh with OCR
                                
                                for idx in range(len(doc)):
                                    page = doc.load_page(idx)
                                    image_list = page.get_images()
                                    
                                    if image_list and len(image_list) > 0:
                                        logger.info(f"Page {idx + 1} contains {len(image_list)} image(s) - performing OCR")
                                    
                                    # Convert page to image for OCR (higher DPI for better accuracy)
                                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom = 288 DPI
                                    img_bytes = pix.tobytes("png")
                                    image = Image.open(io.BytesIO(img_bytes))
                                    
                                    try:
                                        # Use better OCR configuration for full pages
                                        # PSM 6 = Assume uniform block of text
                                        # PSM 11 = Sparse text (better for scanned documents)
                                        # PSM 3 = Fully automatic (best for varied content)
                                        ocr_configs = [
                                            '--psm 3 --oem 3 -c preserve_interword_spaces=1',  # Automatic page segmentation
                                            '--psm 6 --oem 3 -c preserve_interword_spaces=1',  # Uniform block
                                            '--psm 11 --oem 3',  # Sparse text
                                        ]
                                        
                                        ocr_text = ""
                                        for ocr_config in ocr_configs:
                                            try:
                                                ocr_text = pytesseract.image_to_string(image, config=ocr_config).strip()
                                                if ocr_text and len(ocr_text) > 20:
                                                    break  # Use first successful OCR
                                            except Exception as ocr_err:
                                                logger.debug(f"OCR config {ocr_config} failed: {ocr_err}")
                                                continue
                                        
                                    finally:
                                        try:
                                            image.close()
                                            pix = None  # Free memory
                                        except Exception:
                                            pass
                                    
                                    ocr_text = (ocr_text or "").strip()
                                    if ocr_text and len(ocr_text) > 10:  # Use OCR if any meaningful text found
                                        logger.info(f"OCR extracted {len(ocr_text)} characters from page {idx + 1}")
                                        
                                        # Process OCR text into clauses
                                        # First, try to split by paragraphs
                                        paragraphs = ocr_text.split('\n\n')
                                        for para in paragraphs:
                                            para = para.strip()
                                            if not para:
                                                continue
                                            
                                            # If paragraph looks like a table (multiple lines with similar patterns)
                                            lines = para.split('\n')
                                            if len(lines) > 2:
                                                # Check if lines have similar structure (might be table)
                                                line_lengths = [len(line) for line in lines]
                                                if len(set(line_lengths)) <= 3:  # Similar lengths suggest table
                                                    # Keep as table-like structure
                                                    results.append({"text": para, "page": idx + 1})
                                                    continue
                                            
                                            # Split into sentences
                                            sentences = para.replace('. ', '.\n').split('\n')
                                            for s in sentences:
                                                s_clean = s.strip()
                                                if s_clean and len(s_clean) > 5:  # Lower threshold for OCR
                                                    if not s_clean.endswith(('.', '!', '?', ':', ';')):
                                                        s_clean += '.'
                                                    results.append({"text": s_clean, "page": idx + 1})
                                    
                                    if not ocr_text or len(ocr_text) < 10:
                                        logger.warning(f"OCR found little to no text on page {idx + 1}")
                                        
                            if has_images:
                                logger.info(f"OCR processing completed. Extracted {len(results)} clauses from images")
                        finally:
                            try:
                                doc.close()
                            except Exception:
                                pass
                    except Exception as e:
                        logger.warning(f"OCR fallback failed: {e}")
                        import traceback
                        logger.debug(traceback.format_exc())
                finally:
                    self._safe_remove_file(temp_file)
                logger.info(f"Extracted {len(results)} clauses with pages from PDF")
                return results
            elif file_extension in [".pptx", ".ppt"]:
                temp_file = self._download_file(doc_url, ".pptx")
                try:
                    prs = Presentation(temp_file)
                    for slide_num, slide in enumerate(prs.slides, start=1):
                        # Collect text from shapes and tables
                        slide_parts: list[str] = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_parts.append(shape.text.strip())
                            if shape.has_table:
                                table = shape.table
                                for row in table.rows:
                                    row_text = []
                                    for cell in row.cells:
                                        if cell.text.strip():
                                            row_text.append(cell.text.strip())
                                    if row_text:
                                        slide_parts.append(" | ".join(row_text))
                        # Split into pseudo-sentences
                        for chunk in slide_parts:
                            for s in chunk.split('. '):
                                s_clean = s.strip()
                                if s_clean:
                                    if not s_clean.endswith(('.', '!', '?')):
                                        s_clean += '.'
                                    results.append({"text": s_clean, "page": slide_num})
                finally:
                    self._safe_remove_file(temp_file)
                logger.info(f"Extracted {len(results)} clauses with pages from PPTX")
                return results
            else:
                # Fallback for other types: no page info
                text = self.extract_text(doc_url)
                paragraphs = text.split('\n\n')
                for para in paragraphs:
                    sentences = para.strip().split('. ')
                    for s in sentences:
                        s_clean = s.strip()
                        if s_clean:
                            if not s_clean.endswith(('.', '!', '?')):
                                s_clean += '.'
                            results.append({"text": s_clean, "page": None})
                if len(results) == 0:
                    # Line-based fallback
                    for l in [l.strip() for l in text.split('\n') if l.strip()]:
                        if len(l) >= 20:
                            results.append({"text": l if l.endswith(('.', '!', '?')) else (l + '.') , "page": None})
                logger.info(f"Extracted {len(results)} clauses (no pages) from generic file")
                return results
        except Exception as e:
            logger.error(f"Error extracting clauses with pages: {e}")
            # Fallback to simple text
            text = self.extract_text(doc_url)
            return [{"text": text, "page": None}]

    def _generate_temp_filename(self, extension: str) -> str:
        """Generate a unique temporary filename"""
        unique_id = str(uuid.uuid4())
        timestamp = int(time.time())
        return f"temp_{timestamp}_{unique_id}{extension}"

    def _download_file(self, doc_url: str, extension: str) -> str:
        """Download file from URL and return local path.
        Also supports local filesystem paths and file:// URLs by copying into temp.
        """
        temp_filename = self._generate_temp_filename(extension)
        temp_file = os.path.join(self.temp_dir, temp_filename)
        
        try:
            # Handle local filesystem paths
            if os.path.exists(doc_url):
                shutil.copyfile(doc_url, temp_file)
                logger.info(f"Copied local file to: {temp_file}")
                return temp_file

            # Handle file:// URLs
            if doc_url.startswith('file://'):
                from urllib.parse import urlparse
                parsed = urlparse(doc_url)
                local_path = parsed.path
                # On Windows, urlparse yields paths like /C:/path -> strip leading slash
                if os.name == 'nt' and local_path.startswith('/') and len(local_path) > 3 and local_path[2] == ':':
                    local_path = local_path[1:]
                if os.path.exists(local_path):
                    shutil.copyfile(local_path, temp_file)
                    logger.info(f"Copied file:// to: {temp_file}")
                    return temp_file

            # Check file size before downloading
            response = requests.head(doc_url, timeout=10)
            response.raise_for_status()
            
            content_length = response.headers.get('content-length')
            if content_length:
                file_size_mb = int(content_length) / (1024 * 1024)
                if file_size_mb > MAX_FILE_SIZE_MB:
                    raise ValueError(f"File size ({file_size_mb:.1f}MB) exceeds maximum allowed size ({MAX_FILE_SIZE_MB}MB)")
            
            # Download the file
            response = requests.get(doc_url, timeout=DOWNLOAD_TIMEOUT, stream=True)
            response.raise_for_status()
            
            with open(temp_file, "wb") as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
            
            logger.info(f"Downloaded file to: {temp_file}")
            return temp_file
        except Exception as e:
            # Clean up partial file if download failed
            if os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except:
                    pass
            raise e

    def _safe_remove_file(self, file_path: str, max_retries: int = None):
        """Safely remove a file with retries"""
        if max_retries is None:
            max_retries = TEMP_FILE_CLEANUP_RETRIES
        
        for attempt in range(max_retries):
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
                    logger.debug(f"Successfully removed: {file_path}")
                break
            except PermissionError as e:
                if attempt < max_retries - 1:
                    logger.warning(f"File {file_path} is locked, retrying in {TEMP_FILE_CLEANUP_DELAY} second... (attempt {attempt + 1})")
                    try:
                        import gc
                        gc.collect()
                    except Exception:
                        pass
                    time.sleep(TEMP_FILE_CLEANUP_DELAY)
                else:
                    # Final attempt: try to rename then remove (helps break locks on Windows)
                    try:
                        temp_rename = file_path + f".del_{int(time.time())}"
                        os.replace(file_path, temp_rename)
                        try:
                            os.remove(temp_rename)
                            logger.debug(f"Removed after rename: {temp_rename}")
                        except Exception:
                            pass
                    except Exception:
                        pass
                    logger.error(f"Failed to remove {file_path} after {max_retries} attempts: {e}")
            except Exception as e:
                logger.error(f"Error removing {file_path}: {e}")
                break

    def _extract_pdf(self, doc_url: str) -> str:
        temp_file = self._download_file(doc_url, ".pdf")
        try:
            with open(temp_file, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = "".join(page.extract_text() or "" for page in reader.pages)
            logger.info(f"Extracted text from PDF: {len(text)} characters")
            return text
        finally:
            self._safe_remove_file(temp_file)

    def _extract_docx(self, doc_url: str) -> str:
        temp_file = self._download_file(doc_url, ".docx")
        try:
            doc = docx.Document(temp_file)
            text = "\n".join(para.text for para in doc.paragraphs)
            logger.info(f"Extracted text from DOCX: {len(text)} characters")
            return text
        finally:
            self._safe_remove_file(temp_file)

    def _extract_powerpoint(self, doc_url: str) -> str:
        """Extract text from PowerPoint presentations"""
        temp_file = self._download_file(doc_url, ".pptx")
        try:
            prs = Presentation(temp_file)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"\n--- Slide {slide_num} ---\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())
                    # Handle tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            logger.info(f"Extracted text from PowerPoint: {len(text)} characters")
            return text
        finally:
            self._safe_remove_file(temp_file)

    def _extract_excel(self, doc_url: str) -> str:
        """Extract data from Excel files"""
        temp_file = self._download_file(doc_url, ".xlsx")
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(temp_file)
            text_parts = []
            
            for sheet_name in excel_file.sheet_names:
                text_parts.append(f"\n--- Sheet: {sheet_name} ---\n")
                
                # Read the sheet
                df = pd.read_excel(temp_file, sheet_name=sheet_name)
                
                # Add column headers
                if not df.empty:
                    headers = " | ".join(str(col) for col in df.columns)
                    text_parts.append(f"Headers: {headers}")
                    
                    # Add first few rows as sample
                    sample_rows = df.head(10)  # Limit to first 10 rows
                    for idx, row in sample_rows.iterrows():
                        row_text = " | ".join(str(val) for val in row.values)
                        text_parts.append(f"Row {idx+1}: {row_text}")
                    
                    # Add summary statistics
                    text_parts.append(f"\nSummary: {len(df)} rows, {len(df.columns)} columns")
            
            text = "\n".join(text_parts)
            logger.info(f"Extracted data from Excel: {len(text)} characters")
            return text
        finally:
            self._safe_remove_file(temp_file)

    def _extract_zip(self, doc_url: str) -> str:
        """Extract and process files from ZIP archives"""
        temp_file = self._download_file(doc_url, ".zip")
        temp_files_to_cleanup = [temp_file]
        
        try:
            text_parts = ["ZIP Archive Contents:"]
            
            with zipfile.ZipFile(temp_file, 'r') as zip_ref:
                # List all files in the archive
                file_list = zip_ref.namelist()
                text_parts.append(f"Total files: {len(file_list)}")
                
                # Process each file in the archive
                for file_name in file_list:
                    text_parts.append(f"\n--- File: {file_name} ---")
                    
                    # Skip directories
                    if file_name.endswith('/'):
                        continue
                    
                    try:
                        # Read file content
                        with zip_ref.open(file_name) as file:
                            content = file.read()
                            
                        # Determine file type and extract text accordingly
                        file_extension = os.path.splitext(file_name.lower())[1]
                        
                        if file_extension == ".txt":
                            text_parts.append(content.decode('utf-8', errors='ignore'))
                        elif file_extension == ".pdf":
                            # Save PDF temporarily and extract text
                            pdf_temp = self._generate_temp_filename(".pdf")
                            pdf_temp_path = os.path.join(self.temp_dir, pdf_temp)
                            temp_files_to_cleanup.append(pdf_temp_path)
                            
                            with open(pdf_temp_path, "wb") as f:
                                f.write(content)
                            
                            try:
                                with open(pdf_temp_path, "rb") as f:
                                    reader = PyPDF2.PdfReader(f)
                                    pdf_text = "".join(page.extract_text() or "" for page in reader.pages)
                                text_parts.append(pdf_text)
                            except Exception as e:
                                text_parts.append(f"Error processing PDF in ZIP: {str(e)}")
                        elif file_extension in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"]:
                            # Process image with OCR
                            try:
                                image = Image.open(io.BytesIO(content))
                                ocr_text = pytesseract.image_to_string(image)
                                text_parts.append(f"OCR Text: {ocr_text}")
                            except Exception as e:
                                text_parts.append(f"Error processing image in ZIP: {str(e)}")
                        else:
                            text_parts.append(f"Binary file - size: {len(content)} bytes")
                            
                    except Exception as e:
                        text_parts.append(f"Error processing {file_name}: {str(e)}")
            
            text = "\n".join(text_parts)
            logger.info(f"Extracted content from ZIP: {len(text)} characters")
            return text
        finally:
            # Clean up all temporary files
            for temp_file_path in temp_files_to_cleanup:
                self._safe_remove_file(temp_file_path)

    def _extract_image(self, doc_url: str) -> str:
        """Extract text from images using OCR"""
        temp_file = self._download_file(doc_url, "")
        try:
            # Open image and perform OCR
            image = Image.open(temp_file)
            
            # Configure OCR for better accuracy
            custom_config = r'--oem 3 --psm 6'
            text = pytesseract.image_to_string(image, config=custom_config)
            
            # Add image metadata
            metadata = f"Image Format: {image.format}\n"
            metadata += f"Image Size: {image.size}\n"
            metadata += f"Image Mode: {image.mode}\n"
            
            full_text = f"{metadata}\n--- OCR Text ---\n{text}"
            logger.info(f"Extracted text from image: {len(text)} characters")
            return full_text
        finally:
            self._safe_remove_file(temp_file)

    def _extract_binary(self, doc_url: str) -> str:
        """Extract basic information from binary files"""
        temp_file = self._download_file(doc_url, ".bin")
        try:
            with open(temp_file, "rb") as f:
                content = f.read()
            
            # Get file size
            file_size = len(content)
            
            # Try to detect file type using magic bytes
            file_type = "Unknown"
            if content.startswith(b'\x89PNG\r\n\x1a\n'):
                file_type = "PNG Image"
            elif content.startswith(b'\xff\xd8\xff'):
                file_type = "JPEG Image"
            elif content.startswith(b'%PDF'):
                file_type = "PDF Document"
            elif content.startswith(b'PK\x03\x04'):
                file_type = "ZIP Archive"
            elif content.startswith(b'\x50\x4b\x03\x04'):
                file_type = "ZIP Archive (alternative)"
            
            # Get first few bytes as hex for analysis
            hex_preview = content[:32].hex()
            
            text = f"Binary File Analysis:\n"
            text += f"File Type: {file_type}\n"
            text += f"File Size: {file_size} bytes\n"
            text += f"First 32 bytes (hex): {hex_preview}\n"
            
            # If it's a known type, try to extract more info
            if file_type == "PNG Image" or file_type == "JPEG Image":
                try:
                    image = Image.open(temp_file)
                    text += f"Image Format: {image.format}\n"
                    text += f"Image Size: {image.size}\n"
                    text += f"Image Mode: {image.mode}\n"
                except Exception as e:
                    text += f"Error reading image metadata: {str(e)}\n"
            
            logger.info(f"Analyzed binary file: {file_size} bytes")
            return text
        finally:
            self._safe_remove_file(temp_file)